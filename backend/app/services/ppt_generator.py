"""
PPT生成模块
将Markdown格式的大纲转换为PPT文件
"""

import re
from typing import List, Tuple, Optional
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN


@dataclass
class SlideContent:
    """单张幻灯片内容"""
    title: str
    bullets: List[str]
    notes: str


def parse_markdown_outline(markdown_text: str) -> List[SlideContent]:
    """
    解析Markdown格式的大纲，返回结构化的幻灯片列表

    Args:
        markdown_text: Markdown格式的大纲文本

    Returns:
        SlideContent对象列表
    """
    slides: List[SlideContent] = []

    # 按"---"分割幻灯片
    slide_texts = re.split(r'\n---\n', markdown_text)

    for slide_text in slide_texts:
        slide_text = slide_text.strip()
        if not slide_text:
            continue

        # 解析标题
        title_match = re.search(r'^#\s+(?:Slide\s+\d+:\s*)?(.+)$', slide_text, re.MULTILINE)
        title = title_match.group(1).strip() if title_match else "无标题"

        # 解析Notes
        notes_match = re.search(r'^##\s*Notes\s*\n(.+)$', slide_text, re.DOTALL | re.MULTILINE)
        notes = notes_match.group(1).strip() if notes_match else ""

        # 解析要点（排除Notes部分）
        content_without_notes = re.sub(r'^##\s*Notes.*$', '', slide_text, flags=re.DOTALL | re.MULTILINE)

        bullets: List[str] = []
        for line in content_without_notes.split('\n'):
            line = line.strip()
            # 匹配 - 开头的要点
            bullet_match = re.match(r'^-\s+(.+)$', line)
            if bullet_match:
                bullets.append(bullet_match.group(1).strip())

        slides.append(SlideContent(
            title=title,
            bullets=bullets,
            notes=notes
        ))

    return slides


def markdown_to_ppt(markdown_text: str, output_path: str) -> str:
    """
    将Markdown大纲转换为PPT文件

    Args:
        markdown_text: Markdown格式的大纲
        output_path: 输出PPT文件路径

    Returns:
        生成的PPT文件路径
    """
    slides = parse_markdown_outline(markdown_text)

    # 创建演示文稿（16:9比例）
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_content in slides:
        # 使用空白布局
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        if slide_content.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_content.title
            p.font.size = Pt(36)
            p.font.bold = True

        # 添加要点
        if slide_content.bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.6), Inches(11.8), Inches(5.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(slide_content.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "• " + bullet
                p.font.size = Pt(24)
                p.space_after = Pt(14)
                p.level = 0

        # 添加备注
        if slide_content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes

    prs.save(output_path)
    return output_path


def get_slide_count(markdown_text: str) -> int:
    """获取大纲中的幻灯片数量"""
    slides = parse_markdown_outline(markdown_text)
    return len(slides)
